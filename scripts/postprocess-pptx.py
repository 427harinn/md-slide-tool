#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
import json
import unicodedata
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


# ----------------------------
# 基本ユーティリティ
# ----------------------------

def normalize_text(text: str) -> str:
    return unicodedata.normalize("NFKC", (text or "").strip())


def load_images_json(json_path: Path) -> List[dict]:
    with json_path.open("r", encoding="utf-8") as f:
        data = json.load(f)

    if not isinstance(data, list):
        raise ValueError("images json は配列形式である必要があります")

    return data


def emu(value_inch: float) -> int:
    return int(Inches(value_inch))


def rect_right(rect: Tuple[int, int, int, int]) -> int:
    return rect[0] + rect[2]


def rect_bottom(rect: Tuple[int, int, int, int]) -> int:
    return rect[1] + rect[3]


def rect_intersects(a: Tuple[int, int, int, int], b: Tuple[int, int, int, int]) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (
        ax + aw <= bx or
        bx + bw <= ax or
        ay + ah <= by or
        by + bh <= ay
    )


# ----------------------------
# title / body 判定
# ----------------------------

def get_slide_title(slide) -> str:
    placeholder_candidates = []
    fallback_candidates = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        text = normalize_text(shape.text if shape.text else "")
        if not text:
            continue

        top = int(shape.top)
        width = int(shape.width)

        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = int(shape.placeholder_format.type)
            except Exception:
                ph_type = None

            if ph_type in (1, 3):
                score = (1000000000 - top) + width
                placeholder_candidates.append((score, text))
                continue

        score = (1000000000 - top) + width
        fallback_candidates.append((score, text))

    if placeholder_candidates:
        placeholder_candidates.sort(reverse=True)
        return placeholder_candidates[0][1]

    if fallback_candidates:
        fallback_candidates.sort(reverse=True)
        return fallback_candidates[0][1]

    return ""


def find_title_shape(slide):
    placeholder_candidates = []
    fallback_candidates = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        text = normalize_text(shape.text if shape.text else "")
        if not text:
            continue

        top = int(shape.top)
        width = int(shape.width)

        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = int(shape.placeholder_format.type)
            except Exception:
                ph_type = None

            if ph_type in (1, 3):
                score = (1000000000 - top) + width
                placeholder_candidates.append((score, shape))
                continue

        score = (1000000000 - top) + width
        fallback_candidates.append((score, shape))

    if placeholder_candidates:
        placeholder_candidates.sort(reverse=True)
        return placeholder_candidates[0][1]

    if fallback_candidates:
        fallback_candidates.sort(reverse=True)
        return fallback_candidates[0][1]

    return None


def find_body_shape(slide, slide_h: int):
    placeholder_candidates = []
    fallback_candidates = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        width = int(shape.width)
        height = int(shape.height)
        left = int(shape.left)
        top = int(shape.top)
        area = width * height
        text = normalize_text(shape.text if shape.text else "")

        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = int(shape.placeholder_format.type)
            except Exception:
                ph_type = None

            if ph_type in (2, 7):
                score = area
                if top < slide_h * 0.18:
                    score *= 0.9
                placeholder_candidates.append((score, shape, left, top, width, height, text))
                continue

        score = area
        if top < slide_h * 0.22:
            score *= 0.35
        if width < emu(2.5):
            score *= 0.6
        fallback_candidates.append((score, shape, left, top, width, height, text))

    if placeholder_candidates:
        placeholder_candidates.sort(key=lambda x: x[0], reverse=True)
        _, shape, left, top, width, height, _ = placeholder_candidates[0]
        return shape, left, top, width, height

    if fallback_candidates:
        fallback_candidates.sort(key=lambda x: x[0], reverse=True)
        _, shape, left, top, width, height, _ = fallback_candidates[0]
        return shape, left, top, width, height

    return None


# ----------------------------
# 画像サイズ
# ----------------------------

def get_image_size_emu(image_path: Path) -> Tuple[int, int]:
    with Image.open(image_path) as img:
        width_px, height_px = img.size
        dpi = img.info.get("dpi", (96, 96))
        dpi_x = dpi[0] if dpi and dpi[0] else 96
        dpi_y = dpi[1] if dpi and dpi[1] else 96

        width_in = width_px / dpi_x
        height_in = height_px / dpi_y

        return int(Inches(width_in)), int(Inches(height_in))


def fit_size(orig_w: int, orig_h: int, max_w: int, max_h: int) -> Tuple[int, int]:
    if orig_w <= 0 or orig_h <= 0:
        return max_w, max_h

    scale = min(max_w / orig_w, max_h / orig_h)
    scale = min(scale, 1.0)

    return max(1, int(orig_w * scale)), max(1, int(orig_h * scale))


# ----------------------------
# テキスト使用範囲の近似
# ----------------------------

def extract_font_size_pt(shape) -> float:
    sizes = []

    try:
        text_frame = shape.text_frame
    except Exception:
        return 20.0

    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font and run.font.size:
                try:
                    sizes.append(float(run.font.size.pt))
                except Exception:
                    pass

        if para.font and para.font.size:
            try:
                sizes.append(float(para.font.size.pt))
            except Exception:
                pass

    if sizes:
        return max(10.0, sum(sizes) / len(sizes))

    return 20.0


def count_visual_width_units(text: str) -> float:
    units = 0.0
    for ch in text:
        if ch.isspace():
            units += 0.35
            continue

        east = unicodedata.east_asian_width(ch)
        if east in ("F", "W", "A"):
            units += 1.0
        else:
            units += 0.55

    return units


def estimate_text_content_bounds(shape) -> Tuple[int, int, int, int]:
    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    text = shape.text if shape.text else ""

    lines = []
    for raw in text.splitlines():
        s = raw.strip()
        if s:
            lines.append(s)

    if not lines:
        return left, top, min(width, emu(2.0)), min(height, emu(0.8))

    font_size_pt = extract_font_size_pt(shape)
    font_size_emu = font_size_pt * 12700.0

    max_units = max(count_visual_width_units(line) for line in lines)
    line_count = len(lines)

    estimated_text_w = int(max_units * font_size_emu * 0.95 + emu(0.5))
    estimated_text_h = int(line_count * font_size_emu * 1.65 + emu(0.25))

    estimated_text_w = min(max(estimated_text_w, emu(2.0)), width)
    estimated_text_h = min(max(estimated_text_h, emu(1.0)), height)

    return left, top, estimated_text_w, estimated_text_h


def get_occupied_rects(slide, slide_w: int, slide_h: int) -> List[Tuple[int, int, int, int]]:
    rects = []

    title_shape = find_title_shape(slide)
    if title_shape is not None:
        rects.append((
            int(title_shape.left),
            int(title_shape.top),
            int(title_shape.width),
            int(title_shape.height),
        ))

    body_info = find_body_shape(slide, slide_h)
    if body_info:
        body_shape, _, _, _, _ = body_info
        rects.append(estimate_text_content_bounds(body_shape))

    return rects


# ----------------------------
# caption
# ----------------------------

def estimate_caption_box_height(caption: str) -> int:
    if not caption:
        return 0

    length = len(caption.strip())
    if length <= 12:
        return emu(0.22)
    if length <= 28:
        return emu(0.34)
    return emu(0.46)


def add_caption_textbox(slide, caption: str, x: int, y: int, w: int, h: int) -> None:
    if not caption.strip():
        return

    textbox = slide.shapes.add_textbox(x, y, w, h)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = caption
    run.font.size = Pt(10)


# ----------------------------
# 候補矩形
# ----------------------------

def make_candidate_rects(
    slide_w: int,
    slide_h: int,
    occupied_rects: List[Tuple[int, int, int, int]],
) -> List[dict]:
    margin = emu(0.25)
    safe_left = emu(0.45)
    safe_top = emu(0.35)
    safe_right = slide_w - emu(0.45)
    safe_bottom = slide_h - emu(0.35)

    if occupied_rects:
        occ_right = max(rect_right(r) for r in occupied_rects)
        occ_bottom = max(rect_bottom(r) for r in occupied_rects)
    else:
        occ_right = int(slide_w * 0.55)
        occ_bottom = int(slide_h * 0.45)

    right_rect = {
        "name": "right",
        "mode": "vertical",
        "rect": (
            occ_right + margin,
            safe_top,
            max(0, safe_right - (occ_right + margin)),
            max(0, safe_bottom - safe_top),
        ),
    }

    bottom_right_rect = {
        "name": "bottom_right",
        "mode": "single",
        "rect": (
            occ_right + margin,
            occ_bottom + margin,
            max(0, safe_right - (occ_right + margin)),
            max(0, safe_bottom - (occ_bottom + margin)),
        ),
    }

    bottom_rect = {
        "name": "bottom",
        "mode": "horizontal",
        "rect": (
            safe_left,
            occ_bottom + margin,
            max(0, safe_right - safe_left),
            max(0, safe_bottom - (occ_bottom + margin)),
        ),
    }

    return [right_rect, bottom_right_rect, bottom_rect]


# ----------------------------
# レイアウト評価
# ----------------------------

def score_layout(
    images: List[dict],
    rect: Tuple[int, int, int, int],
    mode: str,
) -> Optional[dict]:
    x, y, w, h = rect
    if w <= 0 or h <= 0:
        return None

    gap = emu(0.12)
    count = len(images)

    if count == 0:
        return None

    if count == 1 and mode == "single":
        img = images[0]
        caption_h = estimate_caption_box_height(img.get("caption", ""))
        usable_h = h - caption_h - emu(0.06)
        if usable_h <= 0:
            return None

        orig_w, orig_h = get_image_size_emu(img["resolved_path"])
        pic_w, pic_h = fit_size(orig_w, orig_h, w, usable_h)

        if pic_w < emu(1.3) or pic_h < emu(1.0):
            return None

        px = x + (w - pic_w) // 2
        py = y + max(0, (usable_h - pic_h) // 2)

        return {
            "mode": "single",
            "rect": rect,
            "placements": [{
                "img_x": px,
                "img_y": py,
                "img_w": pic_w,
                "img_h": pic_h,
                "caption_x": x,
                "caption_y": py + pic_h + emu(0.04),
                "caption_w": w,
                "caption_h": caption_h,
            }],
            "score": pic_w * pic_h,
        }

    if mode == "vertical":
        total_gap = gap * max(0, count - 1)
        slot_h = (h - total_gap) // count
        if slot_h <= 0:
            return None

        placements = []
        total_area = 0
        cy = y

        for img in images:
            caption_h = estimate_caption_box_height(img.get("caption", ""))
            usable_h = slot_h - caption_h - emu(0.05)
            if usable_h <= 0:
                return None

            orig_w, orig_h = get_image_size_emu(img["resolved_path"])
            pic_w, pic_h = fit_size(orig_w, orig_h, w, usable_h)

            if pic_w < emu(1.2) or pic_h < emu(0.8):
                return None

            px = x + (w - pic_w) // 2
            py = cy + max(0, (usable_h - pic_h) // 2)

            placements.append({
                "img_x": px,
                "img_y": py,
                "img_w": pic_w,
                "img_h": pic_h,
                "caption_x": x,
                "caption_y": py + pic_h + emu(0.04),
                "caption_w": w,
                "caption_h": caption_h,
            })

            total_area += pic_w * pic_h
            cy += slot_h + gap

        return {
            "mode": "vertical",
            "rect": rect,
            "placements": placements,
            "score": total_area,
        }

    if mode == "horizontal":
        total_gap = gap * max(0, count - 1)
        slot_w = (w - total_gap) // count
        if slot_w <= 0:
            return None

        placements = []
        total_area = 0
        cx = x

        for img in images:
            caption_h = estimate_caption_box_height(img.get("caption", ""))
            usable_h = h - caption_h - emu(0.05)
            if usable_h <= 0:
                return None

            orig_w, orig_h = get_image_size_emu(img["resolved_path"])
            pic_w, pic_h = fit_size(orig_w, orig_h, slot_w, usable_h)

            if pic_w < emu(1.2) or pic_h < emu(0.8):
                return None

            px = cx + (slot_w - pic_w) // 2
            py = y + max(0, (usable_h - pic_h) // 2)

            placements.append({
                "img_x": px,
                "img_y": py,
                "img_w": pic_w,
                "img_h": pic_h,
                "caption_x": cx,
                "caption_y": py + pic_h + emu(0.04),
                "caption_w": slot_w,
                "caption_h": caption_h,
            })

            total_area += pic_w * pic_h
            cx += slot_w + gap

        return {
            "mode": "horizontal",
            "rect": rect,
            "placements": placements,
            "score": total_area,
        }

    return None


def choose_best_layout(
    images: List[dict],
    candidate_rects: List[dict],
    occupied_rects: List[Tuple[int, int, int, int]],
) -> Optional[dict]:
    best = None

    for cand in candidate_rects:
        rect = cand["rect"]
        name = cand["name"]
        mode = cand["mode"]

        if any(rect_intersects(rect, occ) for occ in occupied_rects):
            continue

        if mode == "single" and len(images) != 1:
            continue

        layout = score_layout(images, rect, mode)
        if layout is None:
            continue

        layout["candidate_name"] = name

        if best is None or layout["score"] > best["score"]:
            best = layout

    return best


# ----------------------------
# 画像配置
# ----------------------------

def collect_target_images(slide_info: dict, base_dir: Path) -> List[dict]:
    images = slide_info.get("images", [])
    result = []

    for image in images:
        raw_path = str(image.get("path", "")).strip()
        resolved_path = (base_dir / raw_path).resolve()

        if not raw_path:
            continue

        if not resolved_path.exists():
            print(f"  missing image: {raw_path}")
            continue

        result.append({
            "raw_path": raw_path,
            "resolved_path": resolved_path,
            "caption": str(image.get("caption", "")).strip(),
        })

    return result


def add_offslide_images_near_target_slide(slide, images: List[dict], slide_w: int, anchor_top: int) -> None:
    margin = emu(0.18)
    x = slide_w + emu(0.08)
    y = anchor_top
    max_w = emu(2.5)
    max_h = emu(1.8)

    for image in images:
        caption_h = estimate_caption_box_height(image.get("caption", ""))
        usable_h = max_h - caption_h - emu(0.05)

        orig_w, orig_h = get_image_size_emu(image["resolved_path"])
        pic_w, pic_h = fit_size(orig_w, orig_h, max_w, usable_h)

        slide.shapes.add_picture(str(image["resolved_path"]), x, y, width=pic_w, height=pic_h)

        if caption_h > 0:
            add_caption_textbox(
                slide,
                image["caption"],
                x,
                y + pic_h + emu(0.04),
                max_w,
                caption_h,
            )

        y += pic_h + caption_h + margin


def add_images_to_slide(prs: Presentation, slide, slide_info: dict, base_dir: Path) -> None:
    target_images = collect_target_images(slide_info, base_dir)
    if not target_images:
        return

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    occupied_rects = get_occupied_rects(slide, slide_w, slide_h)

    body_info = find_body_shape(slide, slide_h)
    if body_info:
        _, _, body_top, _, _ = body_info
        offslide_anchor_top = body_top
    else:
        offslide_anchor_top = emu(1.4)

    candidate_rects = make_candidate_rects(
        slide_w=slide_w,
        slide_h=slide_h,
        occupied_rects=occupied_rects,
    )

    best_layout = choose_best_layout(
        images=target_images,
        candidate_rects=candidate_rects,
        occupied_rects=occupied_rects,
    )

    if best_layout is None:
        print("  fallback: off-slide (no candidate fits)")
        add_offslide_images_near_target_slide(
            slide=slide,
            images=target_images,
            slide_w=slide_w,
            anchor_top=offslide_anchor_top,
        )
        return

    print(f"  place: {best_layout['candidate_name']} / {best_layout['mode']}")

    for image, placement in zip(target_images, best_layout["placements"]):
        slide.shapes.add_picture(
            str(image["resolved_path"]),
            placement["img_x"],
            placement["img_y"],
            width=placement["img_w"],
            height=placement["img_h"],
        )

        if placement["caption_h"] > 0 and image.get("caption", "").strip():
            add_caption_textbox(
                slide,
                image["caption"],
                placement["caption_x"],
                placement["caption_y"],
                placement["caption_w"],
                placement["caption_h"],
            )


# ----------------------------
# main
# ----------------------------

def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", required=True, help="後処理対象のPPTXファイル")
    parser.add_argument("--images-json", required=True, help="画像情報JSON")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    images_json_path = Path(args.images_json).resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTXが見つかりません: {pptx_path}")

    if not images_json_path.exists():
        raise FileNotFoundError(f"images json が見つかりません: {images_json_path}")

    slide_infos = load_images_json(images_json_path)
    prs = Presentation(str(pptx_path))

    print(f"PPTX slide count : {len(prs.slides)}")
    print(f"JSON slide count : {len(slide_infos)}")
    print(f"Process count    : {min(len(prs.slides), len(slide_infos))}")

    base_dir = images_json_path.parent

    for i, slide_info in enumerate(slide_infos):
        target_title = normalize_text(slide_info.get("title", ""))
        if not target_title:
            print(f"[WARN] slide title empty at json index {i}")
            continue

        matched_slide = None

        for slide in prs.slides:
            ppt_title = normalize_text(get_slide_title(slide))
            if ppt_title == target_title:
                matched_slide = slide
                break

        if matched_slide is None:
            print(f"[WARN] slide not found: {target_title}")
            continue

        print(
            f"[match] slide {i + 1}: "
            f"title='{target_title}' "
            f"images={len(slide_info.get('images', []))}"
        )

        add_images_to_slide(prs, matched_slide, slide_info, base_dir)

    prs.save(str(pptx_path))
    print(f"saved: {pptx_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())